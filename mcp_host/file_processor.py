"""File processing service -- converts uploaded files into LLM-ready text.

Supported formats
-----------------
Documents : PDF, DOCX, DOC, PPTX, XLSX, XLS, CSV, RTF, TXT, MD
Structured : JSON, XML, HTML / HTM
Images    : JPG, JPEG, PNG, GIF, BMP, WEBP, TIFF, ICO
Audio     : MP3, WAV, M4A, WEBM, OGG, FLAC, AAC
Video     : MP4, AVI, MOV, MKV
Code      : PY, JS, TS, JSX, TSX, JAVA, C, CPP, H, CS, GO, RS, RB, PHP,
            SWIFT, KT, SCALA, SH, BASH, PS1, YAML, YML, TOML, INI, SQL...
"""

import csv
import io
import json
import logging
import os
import re
from pathlib import Path
from typing import Optional, Tuple

logger = logging.getLogger(__name__)

# --------------------------------------------------------------------------- #
# Tunables
# --------------------------------------------------------------------------- #
MAX_TEXT_CHARS = 15_000          # chars forwarded to LLM per file
MAX_FILE_BYTES = 25 * 1024 * 1024  # 25 MB hard limit

# --------------------------------------------------------------------------- #
# Optional dependency guards
# --------------------------------------------------------------------------- #
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.info("PyPDF2 not available - PDF processing disabled")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.info("python-docx not available - Word processing disabled")

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.info("python-pptx not available - PowerPoint processing disabled")

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.info("openpyxl not available - Excel processing disabled")

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False
    logger.info("beautifulsoup4 not available - HTML/XML stripping limited")

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False
    logger.info("striprtf not available - RTF processing disabled")

try:
    from PIL import Image        # noqa: F401 - presence check only
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logger.info("Pillow not available - image resizing disabled")

try:
    import moviepy.editor as moviepy_mp
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False
    logger.info("moviepy not available - video processing disabled")


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #

def _truncate(text: str, limit: int = MAX_TEXT_CHARS) -> str:
    """Trim to *limit* chars and append a note when truncated."""
    if len(text) <= limit:
        return text
    return (
        text[:limit]
        + f"\n\n... [truncated - {len(text):,} total chars, showing first {limit:,}]"
    )


def _mime_from_ext(ext: str) -> str:
    return {
        ".jpg": "jpeg", ".jpeg": "jpeg", ".png": "png", ".gif": "gif",
        ".webp": "webp", ".bmp": "bmp", ".tiff": "tiff", ".tif": "tiff",
    }.get(ext.lower(), "jpeg")


# --------------------------------------------------------------------------- #
# FileProcessor
# --------------------------------------------------------------------------- #

class FileProcessor:
    """Converts uploaded files to plain text for LLM injection."""

    # OpenAI vision models - newest / best quality first
    OPENAI_VISION_MODELS = ["gpt-4o", "gpt-4-turbo", "gpt-4-vision-preview"]

    # HuggingFace caption fallbacks
    HF_VISION_MODELS = [
        "Salesforce/blip-image-captioning-large",
        "Salesforce/blip-image-captioning-base",
        "nlpconnect/vit-gpt2-image-captioning",
        "microsoft/git-base-coco",
    ]

    # Source-code extensions - read as plain text, tagged with language
    CODE_EXTENSIONS = {
        ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h",
        ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
        ".sh", ".bash", ".ps1", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".env", ".sql", ".graphql", ".proto", ".r", ".m", ".lua",
    }

    def __init__(self, openai_client=None, voice_service=None):
        self.openai_client = openai_client
        self.voice_service = voice_service

        self.audio_types  = {".mp3", ".wav", ".m4a", ".webm", ".ogg", ".flac", ".aac"}
        self.video_types  = {".mp4", ".avi", ".mov", ".mkv"}
        self.image_types  = {
            ".jpg", ".jpeg", ".png", ".gif", ".bmp",
            ".webp", ".tiff", ".tif", ".ico",
        }

    # =================================================================== #
    # Public entry point
    # =================================================================== #

    async def process_file(
        self,
        file_data: bytes,
        filename: str,
        user_query: Optional[str] = None,
    ) -> Tuple[str, str]:
        """
        Convert an uploaded file to LLM-ready text.

        Args:
            file_data:   Raw bytes.
            filename:    Original filename (drives type detection).
            user_query:  The user's message - passed to vision so image
                         analysis focuses on what the user actually asked.

        Returns:
            (extracted_text, file_type_tag)
        """
        if len(file_data) > MAX_FILE_BYTES:
            mb = len(file_data) / 1048576
            raise ValueError(
                f"File too large ({mb:.1f} MB). Maximum allowed size is 25 MB."
            )

        ext = Path(filename).suffix.lower()
        logger.info(f"Processing '{filename}' ({len(file_data):,} bytes, ext='{ext}')")

        if ext in self.audio_types:
            return await self._process_audio(file_data, filename), "audio"

        if ext in self.video_types:
            return await self._process_video(file_data, filename), "video"

        if ext in self.image_types:
            return await self._process_image(file_data, filename, user_query), "image"

        if ext == ".pdf":
            return self._process_pdf(file_data, filename), "pdf"

        if ext in {".docx", ".doc"}:
            return self._process_docx(file_data, filename), "docx"

        if ext in {".pptx", ".ppt"}:
            return self._process_pptx(file_data, filename), "pptx"

        if ext in {".xlsx", ".xls"}:
            return self._process_xlsx(file_data, filename), "xlsx"

        if ext == ".csv":
            return self._process_csv(file_data, filename), "csv"

        if ext == ".rtf":
            return self._process_rtf(file_data, filename), "rtf"

        if ext == ".json":
            return self._process_json(file_data, filename), "json"

        if ext == ".xml":
            return self._process_xml(file_data, filename), "xml"

        if ext in {".html", ".htm"}:
            return self._process_html(file_data, filename), "html"

        if ext == ".md":
            return self._decode_text(file_data, filename, label="Markdown"), "markdown"

        if ext in {".txt", ".log", ".text"}:
            return self._decode_text(file_data, filename, label="Text file"), "text"

        if ext in self.CODE_EXTENSIONS:
            lang = ext.lstrip(".")
            return self._decode_text(file_data, filename, label=f"Code ({lang.upper()})"), "code"

        raise ValueError(
            f"Unsupported file type: '{ext}'.\n"
            "Supported: PDF, DOCX, PPTX, XLSX, CSV, RTF, TXT, MD, JSON, XML, HTML, "
            "JPG/PNG/GIF/WEBP (images), MP3/WAV/OGG (audio), MP4/MOV (video), "
            "and most code/config files."
        )

    # =================================================================== #
    # Audio
    # =================================================================== #

    async def _process_audio(self, audio_data: bytes, filename: str) -> str:
        if not self.voice_service:
            return f"[Audio uploaded: '{filename}' - voice service not available]"
        try:
            transcript = await self.voice_service.speech_to_text(audio_data, filename)
            if transcript:
                logger.info(f"Audio transcribed: {len(transcript):,} chars")
                return f"[Audio transcription of '{filename}']\n{transcript}"
            return f"[Audio uploaded: '{filename}' - transcription returned empty]"
        except Exception as e:
            logger.error(f"Audio processing failed: {e}")
            return f"[Audio uploaded: '{filename}' - transcription failed: {e}]"

    # =================================================================== #
    # Video
    # =================================================================== #

    async def _process_video(self, video_data: bytes, filename: str) -> str:
        if not MOVIEPY_AVAILABLE:
            return f"[Video uploaded: '{filename}' - moviepy not installed]"
        if not self.voice_service:
            return f"[Video uploaded: '{filename}' - voice service not available]"
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(
                suffix=Path(filename).suffix, delete=False
            ) as tmp:
                tmp.write(video_data)
                video_path = tmp.name

            clip = moviepy_mp.VideoFileClip(video_path)
            audio_path = video_path + ".mp3"
            clip.audio.write_audiofile(audio_path, logger=None)
            clip.close()

            with open(audio_path, "rb") as f:
                audio_bytes = f.read()

            transcript = await self.voice_service.speech_to_text(audio_bytes, "audio.mp3")
            os.remove(video_path)
            os.remove(audio_path)

            logger.info(f"Video transcribed: {len(transcript):,} chars")
            return f"[Video transcription of '{filename}']\n{transcript}"
        except Exception as e:
            logger.error(f"Video processing failed: {e}")
            return f"[Video uploaded: '{filename}' - processing failed: {e}]"

    # =================================================================== #
    # Image (vision)
    # =================================================================== #

    async def _process_image(
        self,
        image_data: bytes,
        filename: str,
        user_query: Optional[str] = None,
    ) -> str:
        import base64

        # Context-aware prompt - focus on what the user actually asked
        if user_query:
            vision_prompt = (
                f'The user asked: "{user_query}"\n\n'
                "Analyze this image and directly answer the user's question. "
                "Also describe relevant visual details (objects, text, colours, layout)."
            )
        else:
            vision_prompt = (
                "Describe this image thoroughly: content, objects, visible text, "
                "colours, spatial layout, and any other notable details."
            )

        # OpenAI vision - try newest model first
        if self.openai_client:
            b64 = base64.b64encode(image_data).decode()
            mime = _mime_from_ext(Path(filename).suffix)
            for model in self.OPENAI_VISION_MODELS:
                try:
                    resp = self.openai_client.chat.completions.create(
                        model=model,
                        messages=[{
                            "role": "user",
                            "content": [
                                {"type": "text", "text": vision_prompt},
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/{mime};base64,{b64}",
                                        "detail": "auto",
                                    },
                                },
                            ],
                        }],
                        max_tokens=800,
                    )
                    desc = resp.choices[0].message.content
                    logger.info(f"Image analyzed ({model}): {len(desc):,} chars")
                    return f"[Image analysis of '{filename}']\n{desc}"
                except Exception as e:
                    logger.warning(f"OpenAI vision {model} failed: {e}")

        # HuggingFace caption fallback
        hf_token = os.environ.get("HUGGINGFACE_API_KEY") or os.environ.get("HF_TOKEN")
        if hf_token:
            import httpx
            for model in self.HF_VISION_MODELS:
                try:
                    url = f"https://router.huggingface.co/hf-inference/models/{model}"
                    async with httpx.AsyncClient(timeout=30.0) as client:
                        r = await client.post(
                            url,
                            headers={"Authorization": f"Bearer {hf_token}"},
                            content=image_data,
                        )
                        r.raise_for_status()
                        result = r.json()
                    caption = (
                        result[0].get("generated_text", str(result))
                        if isinstance(result, list)
                        else str(result)
                    )
                    logger.info(f"Image captioned (HF/{model}): {len(caption):,} chars")
                    return f"[Image caption of '{filename}']\n{caption}"
                except Exception as e:
                    logger.warning(f"HF vision {model} failed: {e}")

        return (
            f"[Image uploaded: '{filename}' - no vision API available. "
            "Set OPENAI_API_KEY or HUGGINGFACE_API_KEY to enable image analysis.]"
        )

    # =================================================================== #
    # PDF
    # =================================================================== #

    def _process_pdf(self, pdf_data: bytes, filename: str) -> str:
        if not PDF_AVAILABLE:
            return f"[PDF '{filename}' - PyPDF2 not installed]"
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(pdf_data))
            parts = []
            for i, page in enumerate(reader.pages):
                text = (page.extract_text() or "").strip()
                if text:
                    parts.append(f"--- Page {i + 1} ---\n{text}")
            full_text = "\n\n".join(parts)
            logger.info(f"PDF '{filename}': {len(reader.pages)} pages, {len(full_text):,} chars")
            return (
                f"[PDF: '{filename}' - {len(reader.pages)} pages]\n"
                + _truncate(full_text)
            )
        except Exception as e:
            logger.error(f"PDF failed: {e}")
            return f"[PDF upload failed: {e}]"

    # =================================================================== #
    # Word (.docx / .doc)
    # =================================================================== #

    def _process_docx(self, data: bytes, filename: str) -> str:
        if not DOCX_AVAILABLE:
            return f"[Word document '{filename}' - python-docx not installed]"
        try:
            doc = DocxDocument(io.BytesIO(data))
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            table_rows = []
            for table in doc.tables:
                for row in table.rows:
                    cells = " | ".join(c.text.strip() for c in row.cells)
                    if cells.strip():
                        table_rows.append(cells)
            combined = paras + (["[Tables]\n" + "\n".join(table_rows)] if table_rows else [])
            full_text = "\n\n".join(combined)
            logger.info(f"DOCX '{filename}': {len(doc.paragraphs)} paragraphs, {len(full_text):,} chars")
            return f"[Word document: '{filename}']\n" + _truncate(full_text)
        except Exception as e:
            logger.error(f"DOCX failed: {e}")
            return f"[Word document upload failed: {e}]"

    # =================================================================== #
    # PowerPoint (.pptx / .ppt)
    # =================================================================== #

    def _process_pptx(self, data: bytes, filename: str) -> str:
        if not PPTX_AVAILABLE:
            return f"[PowerPoint '{filename}' - python-pptx not installed. Run: pip install python-pptx]"
        try:
            prs = PptxPresentation(io.BytesIO(data))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
            full_text = "\n\n".join(slides)
            logger.info(f"PPTX '{filename}': {len(prs.slides)} slides, {len(full_text):,} chars")
            return (
                f"[PowerPoint: '{filename}' - {len(prs.slides)} slides]\n"
                + _truncate(full_text)
            )
        except Exception as e:
            logger.error(f"PPTX failed: {e}")
            return f"[PowerPoint upload failed: {e}]"

    # =================================================================== #
    # Excel (.xlsx / .xls)
    # =================================================================== #

    def _process_xlsx(self, data: bytes, filename: str) -> str:
        if not XLSX_AVAILABLE:
            return f"[Excel '{filename}' - openpyxl not installed. Run: pip install openpyxl]"
        try:
            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append(" | ".join(cells))
                if rows:
                    sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            full_text = "\n\n".join(sheets_text)
            logger.info(f"XLSX '{filename}': {len(wb.sheetnames)} sheets, {len(full_text):,} chars")
            return (
                f"[Excel: '{filename}' - sheets: {', '.join(wb.sheetnames)}]\n"
                + _truncate(full_text)
            )
        except Exception as e:
            logger.error(f"XLSX failed: {e}")
            return f"[Excel upload failed: {e}]"

    # =================================================================== #
    # CSV
    # =================================================================== #

    def _process_csv(self, data: bytes, filename: str) -> str:
        try:
            text = data.decode("utf-8", errors="replace")
            rows = list(csv.reader(io.StringIO(text)))
            if not rows:
                return f"[CSV '{filename}' is empty]"
            header = " | ".join(rows[0])
            data_rows = [" | ".join(r) for r in rows[1:101]]   # first 100 rows
            note = f"\n... ({len(rows) - 1} total rows)" if len(rows) > 102 else ""
            full_text = f"Columns: {header}\n\n" + "\n".join(data_rows) + note
            logger.info(f"CSV '{filename}': {len(rows)} rows x {len(rows[0])} cols")
            return (
                f"[CSV: '{filename}' - {len(rows) - 1} rows x {len(rows[0])} columns]\n"
                + _truncate(full_text)
            )
        except Exception as e:
            logger.error(f"CSV failed: {e}")
            return f"[CSV upload failed: {e}]"

    # =================================================================== #
    # RTF
    # =================================================================== #

    def _process_rtf(self, data: bytes, filename: str) -> str:
        if not RTF_AVAILABLE:
            return f"[RTF '{filename}' - striprtf not installed. Run: pip install striprtf]"
        try:
            text = rtf_to_text(data.decode("utf-8", errors="replace"))
            logger.info(f"RTF '{filename}': {len(text):,} chars")
            return f"[RTF document: '{filename}']\n" + _truncate(text)
        except Exception as e:
            logger.error(f"RTF failed: {e}")
            return f"[RTF upload failed: {e}]"

    # =================================================================== #
    # JSON
    # =================================================================== #

    def _process_json(self, data: bytes, filename: str) -> str:
        try:
            parsed = json.loads(data.decode("utf-8", errors="replace"))
            pretty = json.dumps(parsed, indent=2, ensure_ascii=False)
            logger.info(f"JSON '{filename}': {len(pretty):,} chars")
            return f"[JSON file: '{filename}']\n" + _truncate(pretty)
        except json.JSONDecodeError as e:
            raw = data.decode("utf-8", errors="replace")
            return f"[JSON file: '{filename}' - parse error: {e}]\nRaw content:\n" + _truncate(raw)

    # =================================================================== #
    # XML
    # =================================================================== #

    def _process_xml(self, data: bytes, filename: str) -> str:
        text = data.decode("utf-8", errors="replace")
        if BS4_AVAILABLE:
            try:
                soup = BeautifulSoup(text, "xml")
                clean = soup.get_text(separator="\n", strip=True)
                logger.info(f"XML '{filename}': {len(clean):,} chars")
                return f"[XML file: '{filename}']\n" + _truncate(clean)
            except Exception:
                pass
        clean = re.sub(r"<[^>]+>", " ", text)
        return f"[XML file: '{filename}']\n" + _truncate(clean)

    # =================================================================== #
    # HTML
    # =================================================================== #

    def _process_html(self, data: bytes, filename: str) -> str:
        text = data.decode("utf-8", errors="replace")
        if BS4_AVAILABLE:
            try:
                soup = BeautifulSoup(text, "html.parser")
                for tag in soup(["script", "style", "meta", "link", "noscript"]):
                    tag.decompose()
                clean = soup.get_text(separator="\n", strip=True)
                logger.info(f"HTML '{filename}': {len(clean):,} chars")
                return f"[HTML page: '{filename}']\n" + _truncate(clean)
            except Exception:
                pass
        clean = re.sub(r"<[^>]+>", " ", text)
        return f"[HTML file: '{filename}']\n" + _truncate(clean)

    # =================================================================== #
    # Plain text / Markdown / Code
    # =================================================================== #

    def _decode_text(self, data: bytes, filename: str, label: str = "Text file") -> str:
        try:
            text = data.decode("utf-8", errors="replace")
            logger.info(f"{label} '{filename}': {len(text):,} chars")
            return f"[{label}: '{filename}']\n" + _truncate(text)
        except Exception as e:
            return f"[{label} '{filename}' - decode failed: {e}]"


# --------------------------------------------------------------------------- #
# Singleton
# --------------------------------------------------------------------------- #

file_processor: Optional[FileProcessor] = None


def initialize_file_processor(
    openai_client=None,
    voice_service=None,
) -> FileProcessor:
    """Create (or replace) the global FileProcessor singleton."""
    global file_processor
    file_processor = FileProcessor(openai_client, voice_service)
    logger.info("FileProcessor initialized")
    return file_processor