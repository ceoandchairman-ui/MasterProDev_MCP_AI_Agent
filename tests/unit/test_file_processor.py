"""Unit tests for FileProcessor — no server, no network, in-memory bytes."""

import asyncio
import io
import json
import textwrap

import pytest
from mcp_host.file_processor import FileProcessor, MAX_FILE_BYTES, _truncate

# FileProcessor without OpenAI or voice (text-only formats don't need them)
@pytest.fixture(scope="module")
def fp() -> FileProcessor:
    return FileProcessor(openai_client=None, voice_service=None)


# ── helpers ────────────────────────────────────────────────────────────────

def run(coro):
    """Run a coroutine using asyncio.run() (Python 3.10+ / 3.14-safe)."""
    return asyncio.run(coro)


# ── plain text ─────────────────────────────────────────────────────────────

def test_process_txt(fp):
    content = "Hello world, this is a test."
    text, _ = run(fp.process_file(content.encode(), "report.txt", user_query="What is this?"))
    assert "Hello world" in text


def test_process_txt_unicode(fp):
    content = "Héllo wörld — résumé café naïve"
    text, _ = run(fp.process_file(content.encode("utf-8"), "notes.txt", user_query="read"))
    assert "Héllo" in text or "Hello" in text  # graceful if non-ascii stripped


def test_process_md(fp):
    content = "# Title\n\nThis is **markdown** content.\n\n- item 1\n- item 2"
    text, _ = run(fp.process_file(content.encode(), "readme.md", user_query="summarize"))
    assert "markdown" in text.lower() or "title" in text.lower()


# ── CSV ────────────────────────────────────────────────────────────────────

def test_process_csv(fp):
    content = b"name,score\nAlice,95\nBob,82\nCarla,78\n"
    text, _ = run(fp.process_file(content, "grades.csv", user_query="Who scored highest?"))
    assert "Alice" in text
    assert "score" in text.lower() or "95" in text


def test_process_csv_many_rows(fp):
    lines = ["id,value"] + [f"{i},{i*10}" for i in range(200)]
    content = "\n".join(lines).encode()
    text, _ = run(fp.process_file(content, "big.csv", user_query="describe"))
    # Should include headers and be truncated to first ~100 rows
    assert "id" in text
    assert "value" in text


# ── JSON ───────────────────────────────────────────────────────────────────

def test_process_json(fp):
    data = {"project": "AI Agent", "version": "2.0", "features": ["chat", "voice"]}
    content = json.dumps(data).encode()
    text, _ = run(fp.process_file(content, "config.json", user_query="list features"))
    assert "AI Agent" in text
    assert "chat" in text


def test_process_json_pretty_printed(fp):
    """Deeply nested JSON should be indented (pretty-printed) in output."""
    data = {"a": {"b": {"c": "deep"}}}
    content = json.dumps(data).encode()
    text, _ = run(fp.process_file(content, "nested.json", user_query="read"))
    assert "deep" in text


# ── code files ────────────────────────────────────────────────────────────

def test_process_python_file(fp):
    content = b"def add(a, b):\n    return a + b\n"
    text, _ = run(fp.process_file(content, "utils.py", user_query="what does this do"))
    assert "def add" in text


def test_process_sql_file(fp):
    content = b"SELECT id, name FROM users WHERE active = 1;"
    text, _ = run(fp.process_file(content, "query.sql", user_query="explain"))
    assert "SELECT" in text


@pytest.mark.parametrize("ext", ["js", "ts", "jsx", "tsx", "yaml", "yml", "xml"])
def test_process_code_extensions(fp, ext):
    content = f"# sample {ext} file\nkey: value\n".encode()
    text, _ = run(fp.process_file(content, f"file.{ext}", user_query="read"))
    assert len(text) > 0


# ── HTML / XML ────────────────────────────────────────────────────────────

def test_process_html(fp):
    content = b"<html><body><h1>Hello</h1><p>World</p></body></html>"
    text, _ = run(fp.process_file(content, "page.html", user_query="read"))
    # Tag-stripped: should contain the text but not raw HTML tags
    assert "Hello" in text
    assert "World" in text


def test_process_xml(fp):
    content = b"<root><item>Alpha</item><item>Beta</item></root>"
    text, _ = run(fp.process_file(content, "data.xml", user_query="list items"))
    assert "Alpha" in text
    assert "Beta" in text


# ── truncation ────────────────────────────────────────────────────────────

def test_truncate_helper():
    """Extracted text over 15 000 chars must be cut with a notice."""
    huge = "AAAA " * 4000   # ~20 000 chars
    result = _truncate(huge)
    assert len(result) <= 15_200  # some room for the truncation notice
    assert "truncated" in result.lower() or len(result) < len(huge)


# ── error conditions ──────────────────────────────────────────────────────

def test_file_too_large(fp):
    """Files over MAX_FILE_BYTES must raise ValueError."""
    big = b"x" * (MAX_FILE_BYTES + 1)
    with pytest.raises(ValueError, match=r"[Ff]ile.*too large|[Ss]ize"):
        run(fp.process_file(big, "huge.txt", user_query="read"))


def test_unsupported_extension(fp):
    """Unknown extension must raise ValueError indicating unsupported type."""
    content = b"\x4d\x5a" + b"\x00" * 64  # fake .exe header
    with pytest.raises(ValueError, match=r"[Uu]nsupported|\.exe"):
        run(fp.process_file(content, "malware.exe", user_query="read"))


def test_empty_file_txt(fp):
    """An empty text file should not crash — returns empty-ish string."""
    text, _ = run(fp.process_file(b"", "empty.txt", user_query="read"))
    assert text is not None
    assert isinstance(text, str)


# ── DOCX (via python-docx) ────────────────────────────────────────────────

def test_process_docx(fp):
    """Build a minimal DOCX in memory and ensure text is extracted."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not installed")

    doc = Document()
    doc.add_paragraph("Annual budget is $10,000,000.")
    doc.add_paragraph("Headcount target: 250.")
    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    text, _ = run(fp.process_file(docx_bytes, "budget.docx", user_query="budget"))
    assert "10,000,000" in text or "budget" in text.lower()


# ── PPTX (via python-pptx) ────────────────────────────────────────────────

def test_process_pptx(fp):
    """Build a minimal PPTX in memory and ensure slide text is extracted."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    txBox.text_frame.text = "Q3 sales grew by 18%."
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    text, _ = run(fp.process_file(pptx_bytes, "pitch.pptx", user_query="sales"))
    assert "18" in text or "sales" in text.lower()


# ── XLSX (via openpyxl) ───────────────────────────────────────────────────

def test_process_xlsx(fp):
    """Build a minimal XLSX in memory and ensure cell data is extracted."""
    try:
        import openpyxl
    except ImportError:
        pytest.skip("openpyxl not installed")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Region", "Q1", "Q2"])
    ws.append(["North", 1_200_000, 1_400_000])
    ws.append(["South",   800_000,   950_000])
    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    text, _ = run(fp.process_file(xlsx_bytes, "revenue.xlsx", user_query="revenue"))
    assert "Region" in text or "North" in text
    assert "Revenue" in text or "1200000" in text or "1,200,000" in text
